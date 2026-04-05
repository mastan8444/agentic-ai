from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import requests
import os
from io import BytesIO
import time
import random
from PIL import Image, ImageDraw, ImageFont
import base64

# 🎨 AI-THEMED COLOR PALETTES
THEMES = {
    "ai_neural": {
        "name": "AI Neural Network",
        "title": RGBColor(98, 0, 234),
        "text": RGBColor(255, 255, 255),
        "accent": RGBColor(0, 255, 255),
        "background": RGBColor(10, 10, 20),
    },
    "ai_future": {
        "name": "AI Future Tech",
        "title": RGBColor(0, 255, 255),
        "text": RGBColor(200, 220, 255),
        "accent": RGBColor(255, 0, 255),
        "background": RGBColor(0, 0, 0),
    },
    "ai_data": {
        "name": "AI Data Science",
        "title": RGBColor(255, 68, 68),
        "text": RGBColor(255, 255, 255),
        "accent": RGBColor(0, 255, 128),
        "background": RGBColor(27, 27, 50),
    },
    "ai_glass": {
        "name": "AI Glass Morphism",
        "title": RGBColor(255, 255, 255),
        "text": RGBColor(230, 230, 255),
        "accent": RGBColor(130, 202, 255),
        "background": RGBColor(20, 30, 45),
    },
    "ai_minimal": {
        "name": "AI Minimal White",
        "title": RGBColor(0, 0, 0),
        "text": RGBColor(50, 50, 60),
        "accent": RGBColor(0, 150, 255),
        "background": RGBColor(255, 255, 255),
    },
    "ai_chip": {
        "name": "AI Chip Circuit",
        "title": RGBColor(0, 255, 100),
        "text": RGBColor(200, 255, 200),
        "accent": RGBColor(0, 200, 255),
        "background": RGBColor(15, 25, 35),
    }
}

DEFAULT_THEME = "ai_neural"
AI_PATTERNS = ["🔮", "🧠", "🤖", "⚡", "💡", "🎯", "📊", "🔄", "✨", "🚀"]

# Create images directory if it doesn't exist
IMAGES_DIR = "generated_images"
os.makedirs(IMAGES_DIR, exist_ok=True)

# 🚀 CREATE PPT
def create_presentation(theme="ai_neural"):
    """Create a new presentation with AI theme"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        if isinstance(theme, dict):
            theme_name = theme.get("name", DEFAULT_THEME).lower()
            theme_key = None
            for key in THEMES:
                if THEMES[key]["name"].lower() == theme_name:
                    theme_key = key
                    break
            prs.theme_choice = theme_key if theme_key else DEFAULT_THEME
        else:
            theme_name = str(theme).lower() if theme else DEFAULT_THEME
            prs.theme_choice = theme_name if theme_name in THEMES else DEFAULT_THEME
        
        prs.ai_theme = THEMES[prs.theme_choice]
        return prs
    except Exception as e:
        print(f"Error creating presentation: {e}")
        prs = Presentation()
        prs.theme_choice = DEFAULT_THEME
        prs.ai_theme = THEMES[DEFAULT_THEME]
        return prs

# 🎯 TITLE SLIDE
def add_title_slide(prs, title, subtitle=None):
    """Add title slide with AI theme"""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = getattr(prs, 'ai_theme', THEMES[DEFAULT_THEME])
        
        # Add background
        add_gradient_background(slide, theme["background"], theme["accent"])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
        tf = title_box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(54)
        run.font.bold = True
        run.font.color.rgb = theme["title"]
        
        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
            sub_tf = sub_box.text_frame
            sub_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            sub_run = sub_tf.paragraphs[0].add_run()
            sub_run.text = subtitle
            sub_run.font.size = Pt(24)
            sub_run.font.color.rgb = theme["text"]
        
        return slide
    except Exception as e:
        print(f"Error adding title slide: {e}")
        return add_fallback_title_slide(prs, title, subtitle)

# 🎨 MAIN SLIDE WITH RELIABLE IMAGES
def add_slide(prs, title, content):
    """Add content slide with reliable image generation"""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = getattr(prs, 'ai_theme', THEMES[DEFAULT_THEME])
        
        # Add gradient background
        add_gradient_background(slide, theme["background"], theme["accent"])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(11), Inches(1.2))
        tf = title_box.text_frame
        icon_text = random.choice(AI_PATTERNS)
        run = tf.paragraphs[0].add_run()
        run.text = f"{icon_text}  {title}"
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = theme["title"]
        
        # Content with bullet points
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(6.5), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.clear()
        
        if content:
            points = [p.strip() for p in content.split("\n") if p.strip()][:8]
            for i, point in enumerate(points):
                if i == 0 and len(tf.paragraphs) > 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = f"• {point}"
                p.font.size = Pt(18)
                p.font.color.rgb = theme["text"]
                p.space_after = Pt(10)
        
        # Generate and add image (RELIABLE VERSION)
        add_generated_image(slide, title, theme)
        
        return True
        
    except Exception as e:
        print(f"Error adding slide: {e}")
        return add_fallback_slide(prs, title, content)

# 🖼️ RELIABLE IMAGE GENERATION FUNCTION
def add_generated_image(slide, title, theme):
    """Generate and add a reliable image to slide"""
    try:
        # Image position
        img_left = Inches(7.5)
        img_top = Inches(1.8)
        img_width = Inches(5.2)
        img_height = Inches(4.8)
        
        # Try multiple image sources
        image_added = False
        
        # Method 1: Try Unsplash (most reliable for free images)
        if not image_added:
            try:
                query = title.replace(" ", "+").replace("&", "").replace("?", "")[:30]
                # Use Unsplash's random image API
                url = f"https://source.unsplash.com/featured/800x600?{query}"
                
                response = requests.get(url, timeout=10, headers={
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                })
                
                if response.status_code == 200:
                    img_path = os.path.join(IMAGES_DIR, f"temp_img_{int(time.time())}.jpg")
                    with open(img_path, "wb") as f:
                        f.write(response.content)
                    
                    if os.path.exists(img_path) and os.path.getsize(img_path) > 1000:
                        slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)
                        image_added = True
                        os.remove(img_path)
            except:
                pass
        
        # Method 2: Try Lorem Picsum (always works)
        if not image_added:
            try:
                random_id = random.randint(1, 200)
                url = f"https://picsum.photos/id/{random_id}/800/600"
                
                response = requests.get(url, timeout=10)
                if response.status_code == 200:
                    img_path = os.path.join(IMAGES_DIR, f"lorem_img_{int(time.time())}.jpg")
                    with open(img_path, "wb") as f:
                        f.write(response.content)
                    
                    if os.path.exists(img_path):
                        slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)
                        image_added = True
                        os.remove(img_path)
            except:
                pass
        
        # Method 3: Create AI-themed decorative image (fallback)
        if not image_added:
            create_ai_decorative_image(slide, title, theme, img_left, img_top, img_width, img_height)
            
    except Exception as e:
        print(f"Image generation error: {e}")
        # Final fallback: Add decorative shape
        add_decorative_shape(slide, theme, img_left, img_top, img_width, img_height)

# 🎨 CREATE AI DECORATIVE IMAGE (Programmatically)
def create_ai_decorative_image(slide, title, theme, left, top, width, height):
    """Create a decorative AI-themed image programmatically"""
    try:
        # Create a PIL image
        img = Image.new('RGB', (800, 600), color=(
            theme["background"].red if hasattr(theme["background"], 'red') else 20,
            theme["background"].green if hasattr(theme["background"], 'green') else 30,
            theme["background"].blue if hasattr(theme["background"], 'blue') else 40
        ))
        draw = ImageDraw.Draw(img)
        
        # Draw geometric patterns
        for i in range(10):
            x1 = random.randint(0, 800)
            y1 = random.randint(0, 600)
            x2 = x1 + random.randint(50, 150)
            y2 = y1 + random.randint(50, 150)
            draw.rectangle([x1, y1, x2, y2], outline=(
                theme["accent"].red if hasattr(theme["accent"], 'red') else 100,
                theme["accent"].green if hasattr(theme["accent"], 'green') else 200,
                theme["accent"].blue if hasattr(theme["accent"], 'blue') else 255
            ), width=3)
        
        # Draw circles (neural network style)
        for _ in range(20):
            x = random.randint(0, 800)
            y = random.randint(0, 600)
            r = random.randint(5, 20)
            draw.ellipse([x-r, y-r, x+r, y+r], fill=(
                theme["title"].red if hasattr(theme["title"], 'red') else 255,
                theme["title"].green if hasattr(theme["title"], 'green') else 255,
                theme["title"].blue if hasattr(theme["title"], 'blue') else 255
            ))
        
        # Add text
        draw.text((400, 300), "AI", fill=(255, 255, 255), anchor="mm")
        draw.text((400, 350), title[:20], fill=(200, 200, 200), anchor="mm")
        
        # Save image
        img_path = os.path.join(IMAGES_DIR, f"decorative_{int(time.time())}.png")
        img.save(img_path)
        
        # Add to slide
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)
            os.remove(img_path)
            
    except Exception as e:
        print(f"Error creating decorative image: {e}")
        add_decorative_shape(slide, theme, left, top, width, height)

def add_decorative_shape(slide, theme, left, top, width, height):
    """Add decorative shape as ultimate fallback"""
    try:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme["accent"]
        shape.fill.transparency = 0.85
        shape.line.fill.background()
        
        # Add text to shape
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.paragraphs[0].text = "🤖\nAI Generated"
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.bold = True
    except:
        pass

# 📊 CHART SLIDE
def add_chart_slide(prs, title, chart_data=None):
    """Add chart slide"""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = getattr(prs, 'ai_theme', THEMES[DEFAULT_THEME])
        
        add_gradient_background(slide, theme["background"], theme["accent"])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = title_box.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = f"📊  {title}"
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = theme["title"]
        
        # Chart data
        if chart_data is None:
            data = CategoryChartData()
            data.categories = ["2020", "2021", "2022", "2023", "2024"]
            data.add_series("Growth", (15, 35, 65, 85, 95))
        else:
            data = chart_data
        
        # Add chart
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1.5), Inches(1.8),
            Inches(7), Inches(4.5),
            data
        )
        
        return True
    except Exception as e:
        print(f"Error adding chart slide: {e}")
        return add_slide(prs, title, "• Chart data visualization\n• Growth metrics")

# 🎨 HELPER FUNCTIONS
def add_gradient_background(slide, color1, color2):
    """Add gradient background"""
    try:
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            Inches(13.333), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = color1
        background.line.fill.background()
        background.z_order = -999
    except:
        pass

def add_fallback_title_slide(prs, title, subtitle):
    """Fallback title slide"""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle if subtitle else "AI Generated"
        return slide
    except:
        return prs.slides.add_slide(prs.slide_layouts[0])

def add_fallback_slide(prs, title, content):
    """Fallback content slide"""
    try:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = content if content else "Content will be added"
        return slide
    except:
        return prs.slides.add_slide(prs.slide_layouts[1])

# 💾 SAVE
def save_presentation(prs, topic):
    """Save presentation"""
    try:
        os.makedirs("output", exist_ok=True)
        clean_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).rstrip()
        filename = f"AI_{clean_topic.replace(' ', '_')}_{int(time.time())}.pptx"
        path = os.path.join("output", filename)
        prs.save(path)
        
        if os.path.exists(path):
            return path
        
        alt_path = os.path.join("output", f"AI_Presentation_{int(time.time())}.pptx")
        prs.save(alt_path)
        return alt_path
    except Exception as e:
        print(f"Error saving: {e}")
        fallback_path = f"AI_Presentation_{int(time.time())}.pptx"
        prs.save(fallback_path)
        return fallback_path

# 🧠 AI BRAIN SLIDE
def add_ai_brain_slide(prs, title, content):
    """Add AI brain themed slide"""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme = getattr(prs, 'ai_theme', THEMES[DEFAULT_THEME])
        
        add_gradient_background(slide, RGBColor(30, 30, 50), RGBColor(10, 10, 30))
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = title_box.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = f"🧠  {title}"
        run.font.size = Pt(38)
        run.font.bold = True
        run.font.color.rgb = theme["title"]
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(6), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.clear()
        
        if content:
            points = [p.strip() for p in content.split("\n") if p.strip()][:6]
            for i, point in enumerate(points):
                if i == 0 and len(tf.paragraphs) > 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = f"⚡  {point}"
                p.font.size = Pt(17)
                p.font.color.rgb = theme["text"]
                p.space_after = Pt(8)
        
        return True
    except Exception as e:
        print(f"Error adding AI brain slide: {e}")
        return add_slide(prs, title, content)